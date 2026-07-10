#!/usr/bin/env python3
"""Generate the Bookly interview deck: Thesis / Architecture / Key Decisions / What's Next.

Structured explicitly around the four prompts a review panel asks:
  1. Thesis — core belief about great CX agents, and how the build reflects it
  2. Architecture — how an inquiry flows through orchestration, tools, memory, prompts
  3. Key decisions — 2-3 choices that mattered most, with trade-offs
  4. What you'd do differently — first change with more time / a production context
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "Bookly_Interview_Deck.pptx"

BG_DARK = RGBColor(15, 20, 25)
ACCENT = RGBColor(79, 140, 255)
TEXT = RGBColor(232, 237, 244)
MUTED = RGBColor(139, 156, 179)
GREEN = RGBColor(111, 207, 151)
RED = RGBColor(242, 153, 74)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullets(tf, items, size=16, color=TEXT):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 and not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def kicker(slide, text: str) -> None:
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5), text, size=14, color=ACCENT, bold=True)


def headline(slide, text: str, size: int = 32) -> None:
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(1.1), text, size=size, bold=True)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ------------------------------------------------------------------
    # Slide 1 — Title
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    kicker(slide, "BOOKLY CUSTOMER SUPPORT AGENT")
    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.4),
                "Architecture & Design Decisions", size=40, bold=True)
    tf = add_textbox(slide, Inches(0.8), Inches(2.9), Inches(10.5), Inches(2.6), "", size=18)
    add_bullets(tf, [
        "1.  The thesis — what a great CX agent believes about itself",
        "2.  Architecture — how one inquiry flows through the system",
        "3.  Key decisions — what I chose, what I traded off, why it was worth it",
        "4.  What I'd change — the first thing, given more time",
    ], size=19)
    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
                "Andrew Child  ·  github.com/andrewwchild/booklydemo  ·  booklydemo.streamlit.app",
                size=12, color=MUTED)
    add_notes(slide, (
        "Thanks for having me — I want to walk through this the way you'd actually want to hear it: "
        "what I believe, how the system reflects that belief, the two or three decisions that mattered "
        "most, and what I'd change next. Then I'm happy to go anywhere you want to dig in."
    ))

    # ------------------------------------------------------------------
    # Slide 2 — The Thesis
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    kicker(slide, "1 · THE THESIS")
    headline(slide, "Never guess. Always check.")
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.7),
                "A wrong answer stated with confidence breaks trust faster than an honest pause ever could.",
                size=18, color=MUTED)
    tf = add_textbox(slide, Inches(0.8), Inches(2.9), Inches(11), Inches(3.6), "", size=17)
    add_bullets(tf, [
        "If the agent is missing information, it asks one focused question — it doesn't guess",
        "If the agent needs a fact, it retrieves it from a real system of record — never from memory",
        "Verification always happens before any action with consequences (refunds, resets)",
        "Being helpful and being right are the same requirement, not a trade-off",
    ], size=18)
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7),
                "Everything on the next slides is this belief, implemented as code.",
                size=16, color=ACCENT, bold=True)
    add_notes(slide, (
        "My core belief is that a support agent earns trust one verified fact at a time, and loses it "
        "instantly the moment it states something wrong with confidence. So I built Bookly's agent around "
        "one rule that governs everything else: never guess, always check. Concretely, that means two "
        "things. First, if the agent doesn't have enough information to act — say, no order number — it "
        "asks one clear question instead of assuming. Second, once it does have enough information, the "
        "actual answer always comes from a real system — the order database, the catalog, the policy "
        "doc — never generated from the model's own memory. You'll see this rule show up three separate "
        "times today: in how the agent clarifies before acting, in how it fetches facts instead of "
        "inventing them, and in how it verifies identity before anything risky, like a refund."
    ))

    # ------------------------------------------------------------------
    # Slide 3 — Architecture: the flow of one inquiry
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    kicker(slide, "2 · ARCHITECTURE")
    headline(slide, "How one inquiry flows through the system")
    flow = (
        "Customer  →  Chat UI  →  BooklyAgent (orchestrator)\n"
        "                              │\n"
        "                              ├─  Memory       — history + \"slots\" filled so far\n"
        "                              ├─  System prompt — the clarify-first house rules\n"
        "                              └─  LLM decides: ask a question, or call a tool?\n"
        "                                        │\n"
        "                                        └─  Tool  →  real Bookly data  →  reply"
    )
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.0), flow, size=17, color=MUTED)
    tf = add_textbox(slide, Inches(0.8), Inches(5.15), Inches(11.5), Inches(1.8), "", size=16)
    add_bullets(tf, [
        "Example: \"Where's my order?\" → missing order ID → agent asks → \"ORD-1001\" → "
        "lookup_order → real UPS tracking, in one plain-language reply",
    ], size=17, color=TEXT)
    add_notes(slide, (
        "Let me walk one message through the system, because the flow matters more than any single "
        "component. A customer types 'Where's my order?' That reaches the orchestrator — the piece of "
        "code that runs the whole loop. The orchestrator hands the message, plus the conversation memory "
        "so far, plus a system prompt describing the house rules, to the language model. The model's only "
        "job at that point is to decide: do I have enough information to act, or do I need to ask "
        "something? Here it doesn't have an order ID, so it asks for one — that's the clarify step. The "
        "customer replies 'ORD-1001'. Memory now has that slot filled, so on the next turn the model "
        "decides it has enough to act, and it calls a tool — lookup_order — which is the only thing "
        "allowed to touch the real order data. That tool comes back with a real UPS tracking number and "
        "delivery estimate, and the model turns that into a normal sentence. Four pieces, one loop: "
        "orchestrator runs it, memory remembers what's been collected, the prompt sets the rules, and "
        "tools are the only door to real facts."
    ))

    # ------------------------------------------------------------------
    # Slide 4 — Architecture: the four components
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    kicker(slide, "2 · ARCHITECTURE")
    headline(slide, "Four components, four jobs")
    rows = [
        ("Orchestrator", "Runs the loop: read message → decide next step → call a tool or reply",
         "agent/orchestrator.py — LLM path + a rules-based fallback with no API key"),
        ("Tools", "The only things allowed to touch real data — one job each",
         "8 functions: lookup_order, verify_customer_identity, initiate_refund, check_stock, "
         "research_books, get_policy, send_password_reset, escalate_to_human"),
        ("Memory", "Tracks conversation history plus the \"slots\" collected so far",
         "order_id, reason, customer_email, book_title — filled in across turns, not re-asked"),
        ("Prompts", "The house rules: clarify-first policy, tone, when to use which tool",
         "One system prompt shared by every intent — no per-intent prompt sprawl"),
    ]
    top = Inches(2.05)
    for i, (name, job, detail) in enumerate(rows):
        y = top + Inches(i * 1.15)
        add_textbox(slide, Inches(0.8), y, Inches(2.2), Inches(0.5), name, size=18, bold=True, color=ACCENT)
        add_textbox(slide, Inches(3.15), y, Inches(9.0), Inches(0.5), job, size=15, color=TEXT)
        add_textbox(slide, Inches(3.15), y + Inches(0.42), Inches(9.0), Inches(0.6), detail, size=13, color=MUTED)
    add_notes(slide, (
        "To define these plainly: the orchestrator is the loop itself — maybe 350 lines I can walk through "
        "line by line, no framework hiding what's happening. Tools are the agent's hands — the only code "
        "path that ever touches real order data, inventory, or policy text, which is exactly what makes "
        "hallucination structurally hard, not just discouraged by a prompt. Memory is short-term working "
        "notes — it's what lets a refund conversation collect an order ID, a reason, and an email across "
        "three separate turns without re-asking anything the customer already said. And the prompt is the "
        "house rules — one shared instruction set that tells the model to clarify before acting and to "
        "always prefer a tool over a guess, rather than writing a different prompt for every single "
        "intent."
    ))

    # ------------------------------------------------------------------
    # Slide 5 — Key decision #1
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    kicker(slide, "3 · KEY DECISION 1 OF 3")
    headline(slide, "Facts come from tools — never from the model", size=30)
    tf = add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.2), "Chose", size=18, bold=True, color=GREEN)
    add_bullets(tf, [
        "LLM's only job is deciding what to do — never stating a transactional fact itself",
        "Every factual claim (status, stock, policy) must pass through a tool call",
        "The system prompt explicitly forbids answering these from memory",
    ], size=16)
    tf = add_textbox(slide, Inches(6.6), Inches(2.1), Inches(5.9), Inches(4.2), "Traded off", size=18, bold=True, color=RED)
    add_bullets(tf, [
        "Every new capability needs a hand-written tool + schema",
        "Slower to add breadth than \"just let the model answer\"",
    ], size=16)
    add_textbox(slide, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.9),
                "Worth it because: a wrong refund amount or order status is a trust failure, not a bug — "
                "the dev cost of one more tool is always cheaper than a hallucinated fact.",
                size=16, color=ACCENT, bold=True)
    add_notes(slide, (
        "The first decision I'd defend hardest is splitting the system into a 'brain' and 'hands.' The "
        "model — the brain — only ever decides what to do next. It never gets to answer a factual question "
        "directly. The tools — the hands — are the only code path allowed to touch real data, and there "
        "are eight of them, each doing exactly one job. What I traded off is real: every new thing the "
        "agent can do requires me to write and register a new tool, which is slower than just letting the "
        "model freehand an answer. But I think that trade is obviously worth it here, because the failure "
        "mode on the other side isn't cosmetic — a wrong tracking number or a wrong refund amount is a "
        "real incident, not an awkward sentence. I'd rather spend the extra engineering time than ship a "
        "system that can sound confident and be wrong."
    ))

    # ------------------------------------------------------------------
    # Slide 6 — Key decision #2
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    kicker(slide, "3 · KEY DECISION 2 OF 3")
    headline(slide, "Clarify-first intent handling, with slot memory", size=30)
    tf = add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.2), "Chose", size=18, bold=True, color=GREEN)
    add_bullets(tf, [
        "Agent tracks a \"pending intent\" + the specific slot it's waiting on",
        "Missing info → one targeted question, not a guess or a wall of questions",
        "Mid-flow, new input answers the open question first — it doesn't get re-diagnosed",
    ], size=16)
    tf = add_textbox(slide, Inches(6.6), Inches(2.1), Inches(5.9), Inches(4.2), "Traded off", size=18, bold=True, color=RED)
    add_bullets(tf, [
        "A lightweight state tracker, not a general-purpose dialogue manager",
        "Works cleanly at 8 intents — wouldn't scale as-is to hundreds",
    ], size=16)
    add_textbox(slide, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.9),
                "Worth it because: refunds need order ID → reason → email, in order — a hand-built slot "
                "tracker makes \"what's still missing\" a one-line check, and keeps behavior testable.",
                size=16, color=ACCENT, bold=True)
    add_notes(slide, (
        "The second decision was how to handle a customer who doesn't give you everything up front, which "
        "is the normal case, not the edge case. I gave the agent a small memory of what it's already "
        "asked for and what it's still waiting on. If something's missing, it asks exactly one focused "
        "question instead of guessing or dumping every question at once. And critically, once the agent "
        "is mid-flow — say, it just asked for an email during a refund — the next thing the customer types "
        "gets interpreted as the answer to that question first, before the system tries to re-guess their "
        "intent from scratch. That fixed a real bug I hit early on, where a reply could accidentally get "
        "reinterpreted as a brand-new request and restart the flow. The trade-off is that this is a "
        "purpose-built tracker, not a general dialogue-management framework — that's fine at eight "
        "well-scoped intents, and I'd say so plainly if this needed to grow to hundreds."
    ))

    # ------------------------------------------------------------------
    # Slide 7 — Key decision #3
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    kicker(slide, "3 · KEY DECISION 3 OF 3")
    headline(slide, "Verify-before-act, as two separate gates", size=30)
    tf = add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.6), Inches(4.2), "Chose", size=18, bold=True, color=GREEN)
    add_bullets(tf, [
        "\"Prove who you are\" (identity) and \"is this allowed\" (eligibility) are two separate tool calls",
        "Refund is only attempted after identity passes",
        "Eligibility can still say no afterward — verification isn't approval",
    ], size=16)
    tf = add_textbox(slide, Inches(6.6), Inches(2.1), Inches(5.9), Inches(4.2), "Traded off", size=18, bold=True, color=RED)
    add_bullets(tf, [
        "Real friction — legitimate customers answer 3-4 questions before a refund starts",
        "A \"just do it\" bot would feel faster in the common case",
    ], size=16)
    add_textbox(slide, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.9),
                "Worth it because: refunds are the highest-risk action in the system — two independent "
                "gates with two distinct failure messages beat one fast, ambiguous one.",
                size=16, color=ACCENT, bold=True)
    add_notes(slide, (
        "The third decision is specifically about refunds, because that's the riskiest action in the "
        "system — wrong person, wrong order, wrong amount. I split 'prove who you are' from 'is this "
        "specific request allowed' into two separate checks. First the agent verifies the customer's "
        "email matches the order on file. Only if that passes does it even attempt the refund — and the "
        "refund tool then separately checks whether the order is actually eligible, for example whether "
        "the return window already closed. Those are genuinely different questions with different correct "
        "responses, so I didn't want them collapsed into one fuzzy 'refund denied' message. The honest "
        "trade-off is friction — a real customer answers three or four questions before anything happens, "
        "and a less careful bot would feel snappier. I think that's the right trade for money leaving the "
        "business, and I can show it live: one refund that succeeds, and one where identity passes but the "
        "return window has already expired — two different correct outcomes, not one confusing error."
    ))

    # ------------------------------------------------------------------
    # Slide 8 — What I'd change
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    kicker(slide, "4 · WHAT I'D DO DIFFERENTLY")
    headline(slide, "First change: let the agent doubt itself, not just the customer", size=28)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.3), Inches(0.7),
                "Today, escalation only fires when the customer explicitly asks for a human.",
                size=17, color=MUTED)
    tf = add_textbox(slide, Inches(0.8), Inches(2.85), Inches(11.3), Inches(1.9), "", size=17)
    add_bullets(tf, [
        "That's reactive — it misses cases where the agent itself is stuck (repeated failed lookups, "
        "the same question rephrased twice) but the customer hasn't asked for a person yet",
        "Confidence-based escalation applies my own thesis — \"never guess\" — to the agent's certainty, "
        "not just to its facts",
    ], size=17)
    add_textbox(slide, Inches(0.8), Inches(4.95), Inches(11.3), Inches(0.5),
                "Also on the list, in order:", size=15, color=MUTED, bold=True)
    tf = add_textbox(slide, Inches(0.8), Inches(5.4), Inches(11.3), Inches(1.6), "", size=15)
    add_bullets(tf, [
        "Real authentication (session/OTP) in place of email-match identity verification",
        "Test the actual AI conversation path, not just the deterministic rules-engine fallback",
        "Swap JSON data files for live OMS/inventory/CRM integrations behind the same tool interface",
    ], size=15, color=TEXT)
    add_notes(slide, (
        "If I had one more week, the first thing I'd change is how the agent decides to escalate. Right "
        "now, handoff to a human only triggers when the customer explicitly asks for one — which is safe, "
        "but it's reactive. It misses the case where the agent itself is struggling — say, a tool keeps "
        "failing, or the customer is rephrasing the same question because nothing's landing — and the "
        "customer hasn't thought to ask for a person yet. I think that's actually the same principle as "
        "everything else I've shown you, just pointed inward: I built the agent to never guess about a "
        "fact, but I haven't yet built it to be honest about its own uncertainty. I didn't ship that here "
        "because a real confidence signal deserves more rigor than I could responsibly build in a "
        "take-home — it's easy to fake, and I didn't want to fake it. After that, in order, I'd replace "
        "email-match verification with real session or OTP-based authentication, extend the eval suite to "
        "actually test the AI conversation path instead of only the deterministic fallback, and swap the "
        "JSON data files for live order-management and inventory integrations — which the tool-based "
        "architecture is specifically designed to make a non-event, since the agent never needs to know "
        "where its data actually comes from."
    ))

    # ------------------------------------------------------------------
    # Slide 9 — Close
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    kicker(slide, "CLOSE")
    headline(slide, "Never guess. Always check. Verify before anything risky.", size=30)
    tf = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.3), Inches(2.6), "", size=18)
    add_bullets(tf, [
        "Thesis → architecture → decisions → roadmap, all pointing at one rule",
        "Depth on core flows beats breadth across shallow ones",
        "Happy to go deeper on any single decision, or run the live demo",
    ], size=19)
    add_notes(slide, (
        "So to bring it back to the start: I believe a support agent's whole job is to be trustworthy on "
        "facts and honest about what it doesn't know yet, and every architectural choice I walked through "
        "today is that belief showing up as code — clarify before acting, get facts from real systems, "
        "verify before anything risky, and be honest with myself about what I'd still change. I'd rather "
        "go deep on a handful of flows I can fully defend than wide across a dozen I can't. Happy to open "
        "up the code, run the live demo, or go deeper on any single decision from here."
    ))

    return prs


if __name__ == "__main__":
    deck = build_deck()
    deck.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
