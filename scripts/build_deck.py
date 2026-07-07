#!/usr/bin/env python3
"""Generate the Bookly CS Agent pitch deck as PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "Bookly_CS_Agent_Pitch.pptx"

BG_DARK = RGBColor(15, 20, 25)
ACCENT = RGBColor(79, 140, 255)
TEXT = RGBColor(232, 237, 244)
MUTED = RGBColor(139, 156, 179)
GREEN = RGBColor(111, 207, 151)
RED = RGBColor(242, 153, 74)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullets(tf, items, size=16, color=TEXT):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 and not tf.paragraphs[0].text else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Title
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6),
                "BOOKLY CUSTOMER SUPPORT AGENT", size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.2),
                "The best CX agents don't guess —\nthey clarify, then act on real data.",
                size=36, bold=True)
    tf = add_textbox(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(3.5), "", size=16)
    add_bullets(tf, [
        "Wrong answers destroy trust faster than one extra question",
        "Human in conversation, machine-precise on facts",
        "Transactional queries always use tools — never free-form generation",
        "One focused clarifying question at a time when info is missing",
    ], size=17)
    add_textbox(slide, Inches(7), Inches(2.8), Inches(5.5), Inches(3.8),
                "User message\n      ↓\nHave required info?\n   ↙     ↘\n Yes      No\n  ↓        ↓\nTool    Clarify",
                size=15, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                "Andrew Child  ·  github.com/andrewwchild/booklydemo  ·  booklydemo.streamlit.app",
                size=12, color=MUTED)
    add_notes(slide, "I'd rather add one turn than ship a confident lie.")

    # Slide 2 — Architecture
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                "Architecture", size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.9),
                "Thin orchestration + LLM reasoning + deterministic tools",
                size=30, bold=True)
    arch = (
        "Customer → Streamlit UI → BooklyAgent (orchestrator)\n"
        "                              ├─ ConversationMemory (slots + history)\n"
        "                              ├─ System prompt (clarify-first rules)\n"
        "                              └─ OpenAI function calling (7 tools)\n"
        "                                    └─ Bookly APIs (orders, catalog, policies)"
    )
    add_textbox(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(2.4), arch, size=15, color=MUTED)
    tf = add_textbox(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(2.3), "", size=16)
    add_bullets(tf, [
        "Orchestrator — LLM ↔ tool loop; rules-based fallback when no API key",
        "Tools — lookup_order · verify_identity · initiate_refund · check_stock · get_policy · send_password_reset · escalate_to_human",
        "Memory — per-session slots (order_id, reason, email, book_title)",
        "Eval harness — 30 golden cases, 100% pass rate, CI-ready",
    ], size=16)
    add_notes(slide, "~350 lines of orchestration I can explain line by line. No LangChain.")

    # Slide 3 — Tools over RAG
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                "KEY DECISION #1", size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.9),
                "Tools over RAG for transactions", size=32, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6),
                "Order status, refunds, and inventory use function calling — not vector search.",
                size=18, color=MUTED)
    tf = add_textbox(slide, Inches(0.8), Inches(2.8), Inches(5.2), Inches(3.5), "Pros", size=18, bold=True, color=GREEN)
    add_bullets(tf, [
        "Zero hallucination risk on order IDs, tracking, stock levels",
        "Actions (initiate_refund) are first-class operations",
        "RAG stays available for long-tail FAQ and policy docs",
    ], size=16)
    tf = add_textbox(slide, Inches(6.5), Inches(2.8), Inches(5.5), Inches(3.5), "Trade-off", size=18, bold=True, color=RED)
    add_bullets(tf, [
        "Each new action needs a tool + integration",
        "More eng work than a single knowledge base",
    ], size=16)
    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
                "Wrong order data is a P0 incident. Transactions need systems of record.",
                size=16, color=ACCENT, bold=True)

    # Slide 4 — Clarify-first + slot memory
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                "KEY DECISION #2", size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.9),
                "Clarify-first + slot memory", size=32, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6),
                "LLM handles natural language; ConversationMemory tracks required fields across turns.",
                size=18, color=MUTED)
    tf = add_textbox(slide, Inches(0.8), Inches(2.8), Inches(5.2), Inches(3.5), "Pros", size=18, bold=True, color=GREEN)
    add_bullets(tf, [
        "Refund: order_id → reason → email → verify → initiate_refund",
        "Regex slot extraction catches structured inputs the LLM might miss",
        "Mid-flow slot collection protected from intent overwrite",
    ], size=16)
    tf = add_textbox(slide, Inches(6.5), Inches(2.8), Inches(5.5), Inches(3.5), "Trade-off", size=18, bold=True, color=RED)
    add_bullets(tf, [
        "Not a full finite-state machine",
        "Production adds auth session pre-fill and tighter validation",
    ], size=16)
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0),
                'Example: "Where\'s my order?" → order ID → lookup_order → UPS tracking',
                size=17, color=ACCENT, bold=True)

    # Slide 5 — Safety & escalation
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                "KEY DECISION #3", size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.9),
                "Verify before act · Escalate with context", size=32, bold=True)
    tf = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), "", size=16)
    add_bullets(tf, [
        "verify_customer_identity before every refund",
        "Verification ≠ approval — eligibility checked separately",
        "Wrong email stops flow before refund record created",
        "escalate_to_human packages transcript + ticket ID",
        "Human agent never makes customer repeat themselves",
    ], size=17)
    add_textbox(slide, Inches(6.5), Inches(2.5), Inches(5.8), Inches(3.5),
                "Refund flow\n\n1. Collect slots\n2. Verify identity\n3. Check eligibility\n4. Initiate return\n\nEscalation\n\n→ Summary + ESC-XXXXXX ticket",
                size=15, color=MUTED)
    add_notes(slide, "Production: OTP, logged-in session, confidence-based auto-escalation.")

    # Slide 6 — Live demo flows
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                "LIVE DEMO", size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.9),
                "Seven flows that prove the thesis", size=32, bold=True)
    rows = [
        ("1. Order lookup", '"Where\'s my order?" → ORD-1001 → UPS tracking'),
        ("2. Delayed shipment", "ORD-1006 → weather delay + new ETA"),
        ("3. Inventory", "Fourth Wing (OOS) · Atomic Habits (in stock)"),
        ("4. Policy FAQ", "Return policy → 30-day window via get_policy"),
        ("5. Refund + verify", "ORD-1003 · damaged · carol@example.com → RFD issued"),
        ("6. Edge case", "ORD-1004 → identity OK, return window expired"),
        ("7. Human handoff", "Escalation → ticket + conversation summary"),
    ]
    top = Inches(2.1)
    for i, (label, detail) in enumerate(rows):
        y = top + Inches(i * 0.72)
        add_textbox(slide, Inches(0.8), y, Inches(2.8), Inches(0.65), label, size=15, bold=True, color=ACCENT)
        add_textbox(slide, Inches(3.6), y, Inches(8.8), Inches(0.65), detail, size=14, color=TEXT)
    add_notes(slide, "Run evals/run_evals.py as backup — 30/30 pass.")

    # Slide 7 — Eval harness
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                "QUALITY", size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.9),
                "Golden-set eval harness — 30 cases, 100% pass", size=32, bold=True)
    tf = add_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4), "", size=16)
    add_bullets(tf, [
        "Tests tool selection, not just phrasing",
        "Covers clarify-first, multi-turn refunds, identity checks",
        "Runs against rules engine in <1 second",
        "CI-ready regression safety when prompts change",
        "Next: LLM-as-judge for tone + shadow-mode staging",
    ], size=17)
    add_textbox(slide, Inches(6.5), Inches(2.5), Inches(5.8), Inches(3),
                "$ python evals/run_evals.py\n\n"
                "30 passed, 0 failed\n"
                "Pass rate: 100%",
                size=18, color=GREEN, align=PP_ALIGN.CENTER)
    add_notes(slide, "Enterprise buyers audit tool-selection accuracy, not vibes.")

    # Slide 8 — Roadmap & close
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5),
                "ROADMAP", size=14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.9),
                "What's next for production", size=32, bold=True)
    rows = [
        ("Today", "Production"),
        ("JSON data stores", "OMS + inventory + CRM integrations"),
        ("Email-match identity", "Auth session + OTP"),
        ("Keyword escalation", "Confidence threshold + sentiment"),
        ("In-process sessions", "Redis + cross-device continuity"),
        ("Rules + LLM dual engine", "LLM-primary + param guardrails"),
    ]
    top = Inches(2.2)
    for i, (left, right) in enumerate(rows):
        y = top + Inches(i * 0.72)
        bold = i == 0
        color = MUTED if i == 0 else TEXT
        sz = 15 if i == 0 else 14
        add_textbox(slide, Inches(0.8), y, Inches(5.2), Inches(0.65), left, size=sz, bold=bold, color=color)
        add_textbox(slide, Inches(6.2), y, Inches(6.2), Inches(0.65), right, size=sz, bold=bold, color=color)
    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9),
                "Depth beats breadth — nail core flows before scaling intents.",
                size=18, color=ACCENT, bold=True)
    add_notes(slide, "Happy to go deeper on any decision. Questions?")

    return prs


if __name__ == "__main__":
    deck = build_deck()
    deck.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
